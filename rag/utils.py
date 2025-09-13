from dotenv import load_dotenv
import os
import serpapi
from pptx import Presentation
from docx import Document
from pathlib import Path
from pptx import Presentation
import pdfplumber
from sentence_transformers import SentenceTransformer
from backend.config import (
    GENERATING_MODEL_NAME,
    EMBEDDING_MODEL_NAME,
    PINECONE_HOST_URL,
)
from pinecone import Pinecone, ServerlessSpec


load_dotenv()

def pptx_content_handler(file_path):
    """Extracts text from pptx with only file name as metadata."""
    prs = Presentation(file_path)
    full_text = []
    file_name = Path(file_path).name

    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        if slide_text.strip():
            full_text.append({
                "file_name": file_name,
                "text": slide_text.strip()
            })

    return full_text

def docx_content_handler(file_path):
    """Extracts text from docx with only file name as metadata."""
    doc = Document(file_path)
    full_text = []
    file_name = Path(file_path).name

    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append({
                "file_name": file_name,
                "text": para.text.strip()
            })

    return full_text

def pdf_content_handler(file_path):
    """Extracts text from pdf with only file name as metadata."""
    full_text = []
    file_name = Path(file_path).name

    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                full_text.append({
                    "file_name": file_name,
                    "text": text.strip()
                })

    return full_text
 
def load_embedding_model(model_name = 'all-MiniLM-L6-v2'):
    """Gets emb model"""
    model = SentenceTransformer(EMBEDDING_MODEL_NAME)
    return model

def load_pinecone():
    PINECONE_API_KEY = os.getenv('PINECONE_API_KEY')
    pc = Pinecone(api_key=PINECONE_API_KEY)

    index_name = "dragggg"  

    if index_name not in pc.list_indexes().names():
        pc.create_index(
            name=index_name,
            dimension=384, 
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1")
        )

    index = pc.Index(index_name)
    return pc, index

    
def load_generating_model():
    """Generating model instatination"""
    model = None

def load_prompts():
    PROMPT_1 = """
You are an AI assistant performing context validation. 
The user asked: "{query}".
The available context is: "{context}".

Determine if the context is sufficient to answer the query. 
If sufficient, set "search_mode": false. 
If insufficient, set "search_mode": true and suggest the missing information in "search_query" (5-6 words). 
Provide a concise reason in "detail".

Strictly respond in JSON format like this:
{
  "search_mode": true or false,
  "detail": "Concise explanation of missing context or why context is enough.",
  "search_query": "Short phrase (5-6 words) to retrieve missing context if needed."
}
"""

    PROMPT_2 = """
You are a helpful AI assistant. 
The user asked: "{query}". 
The available context is: "{context}".

Using the provided context, generate a clear, accurate, and concise answer. 
Include the sources from the context that you used. 

Strictly respond in JSON format like this:
{
  "answer": "Your answer here.",
  "references": ["Reference 1 from context", "Reference 2 from context"]
}
"""

    return {
        "validate_context": PROMPT_1,
        "generate_answer": PROMPT_2
    }

def perform_search(query):
    SERP_API_KEY = os.getenv('SERP_API_KEY')
    client = serpapi.Client(api_key=SERP_API_KEY)
    result = client.search(
        q = query,
        engine = "google",
        location = "India",
        hl = 'en',
        gl = "us",
    )
    snippets = []
    for r in result.get("organic_results", [])[:num_results]:
        snippet_text = r.get("snippet")
        link = r.get("link")
        if snippet_text and link:
            snippets.append({
                "content": snippet_text,
                "link": link
            })

    return snippets